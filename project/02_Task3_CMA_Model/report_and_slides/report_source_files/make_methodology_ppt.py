#!/usr/bin/env python3

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent / "vln_task3_task4_methodology_slides.pptx"


COLORS = {
    "ink": RGBColor(23, 32, 51),
    "muted": RGBColor(93, 105, 128),
    "line": RGBColor(217, 222, 235),
    "soft": RGBColor(247, 249, 253),
    "blue": RGBColor(47, 111, 239),
    "green": RGBColor(0, 166, 118),
    "red": RGBColor(211, 63, 73),
    "orange": RGBColor(255, 159, 28),
    "purple": RGBColor(123, 92, 255),
    "white": RGBColor(255, 255, 255),
}


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=COLORS["line"], width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(shape, text, size=16, bold=False, color=None, align=None):
    shape.text = ""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return shape


def add_box(slide, x, y, w, h, title, body, accent, label=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, COLORS["white"])
    set_line(shape)
    shape.adjustments[0] = 0.12

    if label:
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x + 0.12),
            Inches(y + 0.10),
            Inches(0.72),
            Inches(0.24),
        )
        pill.adjustments[0] = 0.2
        set_fill(pill, accent)
        pill.line.fill.background()
        add_text(pill, label, size=7.8, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)

    title_box = slide.shapes.add_textbox(Inches(x + 0.10), Inches(y + 0.38), Inches(w - 0.20), Inches(0.30))
    add_text(title_box, title, size=12.4, bold=True)
    body_box = slide.shapes.add_textbox(Inches(x + 0.10), Inches(y + 0.72), Inches(w - 0.20), Inches(h - 0.74))
    add_text(body_box, body, size=8.4, color=COLORS["muted"])
    return shape


def add_title(slide, tag, title, subtitle, pill):
    tag_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(3.0), Inches(0.28))
    add_text(tag_box, tag.upper(), size=11, bold=True, color=COLORS["blue"])

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.70), Inches(7.0), Inches(0.55))
    add_text(title_box, title, size=28, bold=True)

    sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.22), Inches(8.4), Inches(0.48))
    add_text(sub_box, subtitle, size=13.2, color=COLORS["muted"])

    pill_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(0.54), Inches(2.7), Inches(0.46)
    )
    pill_shape.adjustments[0] = 0.5
    set_fill(pill_shape, COLORS["soft"])
    set_line(pill_shape)
    add_text(pill_shape, pill, size=10.5, bold=True, color=RGBColor(54, 67, 90), align=PP_ALIGN.CENTER)


def add_lane(slide, x, y, w, h, title, accent):
    lane = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    lane.adjustments[0] = 0.08
    set_fill(lane, RGBColor(252, 253, 255))
    set_line(lane)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.20), Inches(0.13), Inches(0.13))
    set_fill(dot, accent)
    dot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(x + 0.36), Inches(y + 0.10), Inches(w - 0.5), Inches(0.34))
    add_text(tb, title, size=13.4, bold=True)


def add_arrow(slide, x, y):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.35), Inches(0.18))
    set_fill(arrow, RGBColor(142, 151, 170))
    arrow.line.fill.background()


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["white"]

    add_title(
        slide,
        "Task 3 Methodology",
        "VLN Policy Architecture",
        "A multimodal imitation-learning policy that grounds instructions in visual observations and predicts navigation actions.",
        "RGB + Language + Memory → Action",
    )

    add_lane(slide, 0.55, 1.82, 5.9, 1.95, "Visual Stream", COLORS["blue"])
    add_lane(slide, 6.70, 1.82, 5.9, 1.95, "Language Stream", COLORS["purple"])

    add_box(slide, 0.78, 2.34, 1.68, 1.12, "RGB", "Egocentric frame at timestep t.", COLORS["blue"], "Input")
    add_arrow(slide, 2.46, 2.82)
    add_box(slide, 2.75, 2.34, 1.68, 1.12, "ResNet-18", "Extracts layout and object cues.", COLORS["blue"], "Encoder")
    add_arrow(slide, 4.43, 2.82)
    add_box(slide, 4.72, 2.34, 1.45, 1.12, "Visual 256-D", "Projected attention query.", COLORS["blue"], "Feature")

    add_box(slide, 6.93, 2.34, 1.68, 1.12, "Instruction", "Natural-language route command.", COLORS["purple"], "Input")
    add_arrow(slide, 8.61, 2.82)
    add_box(slide, 8.90, 2.34, 1.68, 1.12, "DistilBERT", "Contextual route-token features.", COLORS["purple"], "Encoder")
    add_arrow(slide, 10.58, 2.82)
    add_box(slide, 10.87, 2.34, 1.45, 1.12, "Text Tokens", "Keys/values for attention.", COLORS["purple"], "Tokens")

    y = 4.10
    add_box(slide, 0.85, y, 2.0, 1.12, "CMA Fusion", "Visual query attends to instruction tokens.", COLORS["green"], "Fusion")
    add_arrow(slide, 2.88, y + 0.43)
    add_box(slide, 3.20, y, 2.15, 1.12, "Gated Context", "Balances vision and language.", COLORS["green"], "Gate")
    add_arrow(slide, 5.38, y + 0.43)
    add_box(slide, 5.70, y, 2.15, 1.12, "GRU Controller", "Keeps route memory across turns.", COLORS["orange"], "Memory")
    add_arrow(slide, 7.88, y + 0.43)
    add_box(slide, 8.20, y, 2.0, 1.12, "Policy Head", "Predicts stop, forward, left, right.", COLORS["red"], "Action")
    add_arrow(slide, 10.23, y + 0.43)
    add_box(slide, 10.55, y, 1.7, 1.12, "Cross Entropy", "Imitation loss over actions.", COLORS["red"], "Loss")

    notes = [
        ("Attention", "Visual state attends over DistilBERT tokens."),
        ("Memory", "GRU keeps navigation history across turns."),
        ("Evaluation", "Rollout uses learned actions only, no oracle takeover."),
    ]
    x0 = 1.15
    for i, (t, b) in enumerate(notes):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x0 + i * 3.65),
            Inches(5.55),
            Inches(3.25),
            Inches(0.62),
        )
        shape.adjustments[0] = 0.14
        set_fill(shape, RGBColor(248, 251, 255))
        set_line(shape, RGBColor(219, 230, 255))
        add_text(shape, f"{t}: {b}", size=9.4, color=RGBColor(51, 64, 90))

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.93), Inches(12.2), Inches(0.2))
    add_text(footer, "VLN-CE / Habitat     •     Task 3: Architecture + Training Method", size=8.6, color=RGBColor(135, 145, 166), align=PP_ALIGN.CENTER)


def add_study_card(slide, x, y, number, title, body, color):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.85), Inches(1.82))
    card.adjustments[0] = 0.10
    set_fill(card, COLORS["white"])
    set_line(card)
    icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.18), Inches(y + 0.16), Inches(0.42), Inches(0.42))
    icon.adjustments[0] = 0.18
    set_fill(icon, color)
    icon.line.fill.background()
    add_text(icon, str(number), size=15, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    title_box = slide.shapes.add_textbox(Inches(x + 0.72), Inches(y + 0.14), Inches(1.95), Inches(0.28))
    add_text(title_box, title, size=13.4, bold=True)
    body_box = slide.shapes.add_textbox(Inches(x + 0.20), Inches(y + 0.64), Inches(2.45), Inches(1.04))
    add_text(body_box, body, size=8.5, color=COLORS["muted"])


def slide_task4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["white"]

    add_title(
        slide,
        "Task 4 Methodology",
        "Generalization and Ablation Study",
        "We stress-tested the CMA policy by changing environments, wording, data size, and fusion mechanism.",
        "Evaluate → Stress Test → Compare",
    )

    add_study_card(slide, 0.65, 1.80, 1, "Held-Out Layouts", "Checked transfer beyond familiar navigation patterns.", COLORS["blue"])
    add_study_card(slide, 3.70, 1.80, 2, "Paraphrases", "Rewrote instructions with rule-based lexical substitutions.", COLORS["purple"])
    add_study_card(slide, 6.75, 1.80, 3, "Reduced Data", "Compared smaller and larger route-trajectory subsets.", COLORS["orange"])
    add_study_card(slide, 9.80, 1.80, 4, "Fusion Ablation", "Compared CMA with simpler gated vision-language fusion.", COLORS["green"])

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(4.00), Inches(6.75), Inches(2.15))
    panel.adjustments[0] = 0.08
    set_fill(panel, RGBColor(252, 253, 255))
    set_line(panel)
    title = slide.shapes.add_textbox(Inches(0.90), Inches(4.18), Inches(3.6), Inches(0.28))
    add_text(title, "Interpretation of Results", size=15.4, bold=True)
    bullets = [
        "CMA reduced final distance versus gated fusion.",
        "More route trajectories improved rollout behavior.",
        "Paraphrases exposed wording sensitivity.",
        "Held-out layouts were harder because stopping and turns were less calibrated.",
    ]
    y = 4.58
    for b in bullets:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.92), Inches(y + 0.08), Inches(0.08), Inches(0.08))
        set_fill(dot, COLORS["blue"])
        dot.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(1.12), Inches(y), Inches(5.8), Inches(0.32))
        add_text(tb, b, size=9.8, color=RGBColor(48, 59, 81))
        y += 0.38

    metric_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.72), Inches(4.00), Inches(4.85), Inches(2.15))
    metric_panel.adjustments[0] = 0.08
    set_fill(metric_panel, RGBColor(252, 253, 255))
    set_line(metric_panel)
    mt = slide.shapes.add_textbox(Inches(7.96), Inches(4.18), Inches(2.5), Inches(0.28))
    add_text(mt, "Key Numbers", size=15.4, bold=True)
    metrics = [
        ("25%", "Best clean validation SR", COLORS["green"]),
        ("5.92m", "CMA distance in ablation", COLORS["blue"]),
        ("37.5%", "Best controlled SR/SPL", COLORS["orange"]),
        ("↓", "Paraphrase robustness drop", COLORS["red"]),
    ]
    positions = [(7.95, 4.60), (10.15, 4.60), (7.95, 5.34), (10.15, 5.34)]
    for (value, label, color), (x, y) in zip(metrics, positions):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.95), Inches(0.62))
        box.adjustments[0] = 0.15
        set_fill(box, COLORS["white"])
        set_line(box)
        v = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.08), Inches(0.72), Inches(0.28))
        add_text(v, value, size=18.5, bold=True, color=color)
        l = slide.shapes.add_textbox(Inches(x + 0.78), Inches(y + 0.08), Inches(1.08), Inches(0.44))
        add_text(l, label, size=6.8, color=COLORS["muted"])

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.93), Inches(12.2), Inches(0.2))
    add_text(footer, "Task 4: Generalization + Ablations     •     Interpretation-focused evaluation", size=8.6, color=RGBColor(135, 145, 166), align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_architecture(prs)
    slide_task4(prs)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
